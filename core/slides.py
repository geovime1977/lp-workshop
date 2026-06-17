from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import io
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np
from core.fmt import br

# ── Paleta ────────────────────────────────────────────────────────────
VERDE       = RGBColor(0x2D, 0x6A, 0x4F)
VERDE_MED   = RGBColor(0x40, 0x96, 0x6B)
VERDE_CLARO = RGBColor(0x52, 0xB7, 0x88)
AMARELO     = RGBColor(0xF4, 0xD0, 0x3F)
BRANCO      = RGBColor(0xFF, 0xFF, 0xFF)
CINZA_ESCURO= RGBColor(0x1A, 0x2E, 0x1D)
PRETO_VERDE = RGBColor(0x0A, 0x1A, 0x0D)
VERMELHO    = RGBColor(0xE7, 0x4C, 0x3C)

W = Inches(13.33)
H = Inches(7.5)
M = Inches(0.5)
CW = W - 2 * M

CORES_CULT = ["#52B788", "#95D5B2", "#2D6A4F", "#F4D03F"]


# ── Helpers ───────────────────────────────────────────────────────────
def _slide(prs, bg=PRETO_VERDE):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    f = s.background.fill
    f.solid()
    f.fore_color.rgb = bg
    return s


def _rect(slide, l, t, w, h, color):
    sh = slide.shapes.add_shape(1, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh


def _txt(slide, l, t, w, h, text, size=18, bold=False,
         color=BRANCO, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def _header(slide, titulo, subtitulo=""):
    _rect(slide, 0, 0, W, Inches(1.15), VERDE)
    _rect(slide, 0, 0, Inches(0.08), Inches(1.15), AMARELO)
    _txt(slide, Inches(0.2), Inches(0.18), CW, Inches(0.6),
         titulo, size=28, bold=True, color=BRANCO)
    if subtitulo:
        _txt(slide, Inches(0.2), Inches(0.72), CW, Inches(0.38),
             subtitulo, size=14, color=VERDE_CLARO)


def _fig_to_bytes(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=200, bbox_inches="tight",
                facecolor=fig.get_facecolor())
    buf.seek(0)
    return buf


def _add_img(slide, buf, l, t, w, h):
    slide.shapes.add_picture(buf, l, t, w, h)


def _grafico_barras(culturas, qtds, title):
    fig, ax = plt.subplots(figsize=(7, 4.5), facecolor="#F0F4F0")
    ax.set_facecolor("#FFFFFF")
    bars = ax.bar(culturas, qtds, color=["#2D6A4F", "#40966B", "#52B788", "#F4D03F"],
                  edgecolor="#1A4035", linewidth=1.2, width=0.55)
    for bar, q in zip(bars, qtds):
        ax.text(bar.get_x() + bar.get_width() / 2,
                bar.get_height() + max(qtds) * 0.025,
                f"{br(q, 0)} ha", ha="center", va="bottom",
                color="#1A2E1D", fontsize=11, fontweight="bold")
    ax.set_ylabel("Hectares", color="#1A2E1D", fontsize=12, fontweight="bold")
    ax.set_title(title, color="#1A2E1D", fontsize=13, fontweight="bold", pad=12)
    ax.tick_params(colors="#1A2E1D", labelsize=10)
    ax.set_ylim(0, max(qtds) * 1.2)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    for sp in ["bottom", "left"]:
        ax.spines[sp].set_edgecolor("#2D6A4F")
        ax.spines[sp].set_linewidth(1.5)
    ax.yaxis.grid(True, color="#CCDDCC", linewidth=0.7, alpha=0.8)
    ax.set_axisbelow(True)
    plt.tight_layout()
    return fig


def _grafico_pizza(culturas, receitas, lucro):
    fig, ax = plt.subplots(figsize=(5.5, 4.5), facecolor="#F0F4F0")
    wedges, texts, autotexts = ax.pie(
        receitas, labels=culturas, autopct="%1.1f%%",
        colors=["#2D6A4F", "#40966B", "#52B788", "#F4D03F"],
        textprops={"color": "#1A2E1D", "fontsize": 10, "fontweight": "bold"},
        pctdistance=0.78,
        wedgeprops={"edgecolor": "white", "linewidth": 2},
        startangle=140,
    )
    for at in autotexts:
        at.set_fontsize(9)
        at.set_color("#1A2E1D")
        at.set_fontweight("bold")
    ax.set_title(f"Receita Total: R$ {br(lucro, 0)}",
                 color="#1A2E1D", fontsize=11, fontweight="bold", pad=12)
    plt.tight_layout()
    return fig


def _grafico_recursos(uso):
    labels = {"terra": "Terra", "orcamento": "Orçamento",
               "agua": "Água", "mao_obra": "Mão de Obra"}
    keys = list(uso.keys())
    pcts = [uso[k][0] / uso[k][1] * 100 for k in keys]
    nomes = [labels[k] for k in keys]

    fig, ax = plt.subplots(figsize=(7.5, 3.8), facecolor="#F0F4F0")
    ax.set_facecolor("#FFFFFF")
    y = np.arange(len(nomes))
    ax.barh(y, [100]*len(y), color="#E8F0E8", height=0.55, edgecolor="#CCDDCC", linewidth=0.8)
    cores = ["#C0392B" if p >= 99 else "#2D6A4F" for p in pcts]
    ax.barh(y, pcts, color=cores, height=0.55, edgecolor="white", linewidth=0.5)
    for i, (p, k) in enumerate(zip(pcts, keys)):
        u, t, un = uso[k]
        ax.text(102, i, f"{p:.1f}%  ({br(u, 0)} / {br(t, 0)} {un})",
                va="center", color="#1A2E1D", fontsize=9, fontweight="bold")
    ax.set_yticks(y)
    ax.set_yticklabels(nomes, color="#1A2E1D", fontsize=11, fontweight="bold")
    ax.set_xlim(0, 175)
    ax.set_xlabel("Utilização (%)", color="#1A2E1D", fontsize=10)
    ax.tick_params(colors="#1A2E1D", labelsize=9)
    ax.set_title("Uso dos Recursos — Gargalos em Vermelho",
                 color="#1A2E1D", fontsize=12, fontweight="bold", pad=10)
    for sp in ax.spines.values():
        sp.set_visible(False)
    ax.xaxis.grid(True, color="#CCDDCC", linewidth=0.7)
    ax.set_axisbelow(True)
    ativo = mpatches.Patch(color="#C0392B", label="Gargalo (100% consumido)")
    folga = mpatches.Patch(color="#2D6A4F", label="Com folga")
    ax.legend(handles=[ativo, folga], loc="lower right",
              facecolor="#F0F4F0", edgecolor="#CCDDCC", fontsize=9)
    plt.tight_layout()
    return fig


# ── Gerador principal ─────────────────────────────────────────────────
def gerar_pptx(resultado: dict) -> bytes:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    culturas = resultado["culturas"]
    qtds     = resultado["quantidades"]
    receitas = resultado["receita_por_cultura"]
    lucro    = resultado["lucro_otimo"]
    uso      = resultado["uso"]
    params   = resultado.get("params", {
        "margem":   [3500, 2300, 4200, 3400],
        "custo":    [4500, 3200, 6000, 2800],
        "agua":     [600,  800,  500,  1500],
        "mao_obra": [12,   15,   22,   8],
        "demanda_max": [2500, 2000, 1500, 1000],
    })

    # ── S1 Capa ────────────────────────────────────────────────────────
    s = _slide(prs, VERDE)
    _rect(s, 0, Inches(2.6), W, Inches(2.5), PRETO_VERDE)
    _rect(s, 0, Inches(5.1), W, Inches(0.05), AMARELO)
    _txt(s, M, Inches(0.35), CW, Inches(0.6),
         "MBA EM PESQUISA OPERACIONAL E TOMADA DE DECISÃO — BSBr",
         size=14, color=AMARELO, align=PP_ALIGN.CENTER)
    _txt(s, M, Inches(0.95), CW, Inches(1.4),
         "Otimização do Mix de Culturas", size=42, bold=True,
         color=BRANCO, align=PP_ALIGN.CENTER)
    _txt(s, M, Inches(2.4), CW, Inches(0.5),
         "Programação Linear Aplicada ao Agronegócio Brasileiro",
         size=20, color=AMARELO, align=PP_ALIGN.CENTER)
    _txt(s, M, Inches(2.85), CW, Inches(0.6),
         "Cooperativa AgroPrime — Caso Hipotético Didático", size=17,
         color=BRANCO, align=PP_ALIGN.CENTER)
    _txt(s, M, Inches(3.5), CW, Inches(0.5),
         "Ferramenta: Python + PuLP (CBC Solver)",
         size=14, color=VERDE_CLARO, align=PP_ALIGN.CENTER, italic=True)
    _txt(s, M, Inches(6.6), CW, Inches(0.5),
         "Prof. Gabriel Capela   |   Workshop — Jun/2026   |   Etapas I, II e III",
         size=13, color=BRANCO, align=PP_ALIGN.CENTER)

    # ── S2 Sobre a Ferramenta ─────────────────────────────────────────
    s = _slide(prs)
    _header(s, "AgroPrime LP — Sobre a Ferramenta",
            "App Streamlit · MBA em PO · Prof. Gabriel Capela — BSBr · Entrega: 20/06/2026")

    col1_w = Inches(6.2)
    col2_x = Inches(7.0)
    col2_w = W - col2_x - M

    # -- coluna esquerda --
    _txt(s, M, Inches(1.28), col1_w, Inches(0.34),
         "Descrição", size=14, bold=True, color=AMARELO)
    _txt(s, M, Inches(1.62), col1_w, Inches(0.75),
         "App Streamlit de Programação Linear para o workshop do MBA em "
         "Pesquisa Operacional (Prof. Gabriel Capela — BSBr).",
         size=12, color=BRANCO)

    _txt(s, M, Inches(2.42), col1_w, Inches(0.34),
         "Problema", size=14, bold=True, color=AMARELO)
    _txt(s, M, Inches(2.76), col1_w, Inches(1.15),
         "Cooperativa AgroPrime precisa decidir quantos hectares alocar entre Soja, "
         "Milho, Algodão e Cana-de-Açúcar para maximizar margem de contribuição "
         "respeitando restrições de área (5.000 ha), orçamento (R$ 20M), "
         "água (4M m³) e mão de obra (60.000 hh).",
         size=12, color=BRANCO)

    _txt(s, M, Inches(4.0), col1_w, Inches(0.34),
         "Stack", size=14, bold=True, color=AMARELO)
    _txt(s, M, Inches(4.34), col1_w, Inches(0.4),
         "Python · Streamlit · PuLP · NumPy · Pandas · Matplotlib · python-pptx · fpdf2",
         size=11, color=VERDE_CLARO)

    _txt(s, M, Inches(4.85), col1_w, Inches(0.34),
         "Acesso & Localização", size=14, bold=True, color=AMARELO)
    for i, line in enumerate([
        "Repositório: github.com/geovime1977/lp-workshop",
        "Local: /Users/virmecati/projetos/lp-workshop/",
        "Porta: 8505  |  Backup: onedrive-eixoestrategico10:projetos/lp-workshop",
    ]):
        _txt(s, M, Inches(5.19) + Inches(0.3)*i, col1_w, Inches(0.28),
             line, size=10, color=VERDE_CLARO)

    _txt(s, M, Inches(6.15), col1_w, Inches(0.34),
         "Como rodar", size=14, bold=True, color=AMARELO)
    _rect(s, M, Inches(6.5), col1_w, Inches(0.7), RGBColor(0x0D, 0x1B, 0x10))
    _txt(s, M + Inches(0.1), Inches(6.55), col1_w - Inches(0.2), Inches(0.6),
         "cd ~/projetos/lp-workshop\n"
         ".venv/bin/streamlit run app.py --server.port 8505\n"
         "# http://localhost:8505",
         size=9, color=VERDE_CLARO)

    # -- coluna direita: módulos --
    _rect(s, col2_x, Inches(1.25), col2_w, Inches(0.5), VERDE)
    _txt(s, col2_x + Inches(0.1), Inches(1.3), col2_w, Inches(0.45),
         "Módulos (6 páginas)", size=14, bold=True, color=AMARELO)

    modulos = [
        ("Tutorial",             "Guia de uso + conceitos LP"),
        ("Contextualização",     "Dados da AgroPrime"),
        ("Modelagem Matemática", "Função objetivo + restrições"),
        ("Solver Interativo",    "PuLP + análise de sensibilidade"),
        ("Resultados e Slides",  "PPTX 12 slides + PDF modelagem"),
        ("Método Gráfico",       "Visualização 2D região viável"),
    ]
    bgs_mod = [CINZA_ESCURO, RGBColor(0x0F, 0x22, 0x14)] * 3
    for i, (nome, desc) in enumerate(modulos):
        top = Inches(1.78) + Inches(0.96)*i
        _rect(s, col2_x, top, col2_w, Inches(0.9), bgs_mod[i])
        _txt(s, col2_x + Inches(0.1), top + Inches(0.05),
             col2_w - Inches(0.2), Inches(0.38),
             nome, size=12, bold=True, color=VERDE_CLARO)
        _txt(s, col2_x + Inches(0.1), top + Inches(0.46),
             col2_w - Inches(0.2), Inches(0.38),
             desc, size=11, color=BRANCO, italic=True)

    # ── S3 Pipeline ────────────────────────────────────────────────────
    s = _slide(prs)
    _header(s, "Pipeline da Solução",
            "Fluxo de processamento — da interface ao entregável final")

    stages = [
        ("INTERFACE",    "Streamlit",         VERDE,       BRANCO,      "Coleta parâmetros\ndo usuário"),
        ("INTELIGÊNCIA", "PuLP (Solver LP)",  VERDE_MED,   BRANCO,      "Resolve o modelo LP\nvia Simplex / CBC"),
        ("SUPORTE",      "Pandas & NumPy",    VERDE_CLARO, PRETO_VERDE, "Processa e formata\nos resultados"),
        ("ENTREGA",      "PDF & PPTX",        AMARELO,     PRETO_VERDE, "Gera os relatórios\nexportáveis"),
    ]

    box_w   = Inches(2.3)
    box_h   = Inches(2.7)
    arr_w   = Inches(0.75)
    total_w = 4 * box_w + 3 * arr_w
    start_x = (W - total_w) / 2
    box_top = Inches(2.1)

    for i, (stage, tech, bg, fg, desc) in enumerate(stages):
        bx = start_x + i * (box_w + arr_w)

        _rect(s, bx, box_top, box_w, box_h, bg)

        _txt(s, bx, box_top + Inches(0.18), box_w, Inches(0.48),
             stage, size=15, bold=True, color=fg, align=PP_ALIGN.CENTER)

        div_color = BRANCO if i < 2 else CINZA_ESCURO
        _rect(s, bx + Inches(0.2), box_top + Inches(0.74),
              box_w - Inches(0.4), Inches(0.03), div_color)

        _txt(s, bx, box_top + Inches(0.85), box_w, Inches(0.55),
             tech, size=14, bold=True, color=fg, align=PP_ALIGN.CENTER)

        _txt(s, bx, box_top + Inches(1.5), box_w, Inches(0.85),
             desc, size=11, color=fg, align=PP_ALIGN.CENTER, italic=True)

        if i < 3:
            ax = bx + box_w
            cy = box_top + box_h / 2
            arr = s.shapes.add_shape(
                13,
                ax + Inches(0.05), cy - Inches(0.22),
                arr_w - Inches(0.1), Inches(0.44),
            )
            arr.fill.solid()
            arr.fill.fore_color.rgb = AMARELO
            arr.line.fill.background()

    _txt(s, M, Inches(5.1), CW, Inches(0.4),
         "Cada camada é independente: o solver pode ser trocado sem alterar a interface ou os entregáveis.",
         size=13, color=VERDE_CLARO, align=PP_ALIGN.CENTER, italic=True)

    # ── S4 Agenda ──────────────────────────────────────────────────────
    s = _slide(prs)
    _header(s, "Agenda")
    itens = [
        ("I",   "Contextualização",      "Empresa, problema e objetivo de negócio"),
        ("II",  "Modelagem Matemática",  "Variáveis, função objetivo e restrições"),
        ("III", "Solução Computacional", "Python / PuLP — código e resultados"),
        ("IV",  "Resolução Manual",      "Sistema 2×2 — identificar ativas e resolver por substituição"),
        ("V",   "Cálculos e Verificação","Prova da solução ótima passo a passo"),
        ("VI",  "Visualizações",         "Gráficos de alocação, receita e uso de recursos"),
        ("VII", "Conclusão",             "Insights e aplicabilidade para cooperativas brasileiras"),
    ]
    for i, (num, titulo, desc) in enumerate(itens):
        top = Inches(1.3) + Inches(0.74) * i
        _rect(s, M, top, Inches(0.6), Inches(0.55), VERDE)
        _txt(s, M + Inches(0.04), top + Inches(0.06), Inches(0.52), Inches(0.45),
             num, size=13, bold=True, color=AMARELO, align=PP_ALIGN.CENTER)
        _txt(s, Inches(1.2), top + Inches(0.01), Inches(4.5), Inches(0.32),
             titulo, size=15, bold=True, color=BRANCO)
        _txt(s, Inches(1.2), top + Inches(0.30), Inches(11.5), Inches(0.32),
             desc, size=12, color=VERDE_CLARO)

    # ── S3 Contextualização ────────────────────────────────────────────
    s = _slide(prs)
    _header(s, "Etapa I — Contextualização do Problema",
            "AgroPrime · Empresa Fictícia · Caso Didático · Safra 2025/2026")

    col1_w = Inches(6.0)
    col2_x = Inches(7.0)
    col2_w = W - col2_x - M

    _txt(s, M, Inches(1.3), col1_w, Inches(0.4),
         "A Empresa", size=16, bold=True, color=AMARELO)
    _txt(s, M, Inches(1.7), col1_w, Inches(1.4),
         "A Cooperativa AgroPrime é uma empresa fictícia criada para fins didáticos, "
         "inspirada no perfil típico de cooperativas do Centro-Oeste brasileiro. "
         "Com 5.000 hectares disponíveis para a safra 2025/2026, a diretoria "
         "precisa definir o mix ótimo de culturas.", size=14, color=BRANCO)
    _txt(s, M, Inches(3.1), col1_w, Inches(0.4),
         "A Decisão", size=16, bold=True, color=AMARELO)
    _txt(s, M, Inches(3.5), col1_w, Inches(1.8),
         "Quantos hectares plantar de cada cultura — Soja, Milho, Algodão "
         "e Cana-de-Açúcar — para MAXIMIZAR O LUCRO TOTAL da safra, "
         "respeitando limitações de:\n"
         "  • Área cultivável disponível\n"
         "  • Orçamento de insumos\n"
         "  • Disponibilidade hídrica\n"
         "  • Mão de obra na safra", size=14, color=BRANCO)

    _rect(s, col2_x, Inches(1.2), col2_w, Inches(0.5), VERDE)
    _txt(s, col2_x + Inches(0.1), Inches(1.25), col2_w, Inches(0.45),
         "Recursos Disponíveis", size=14, bold=True, color=AMARELO)
    recursos_info = [
        ("🌍 Terra",       "5.000 ha"),
        ("💰 Orçamento",   "R$ 20.000.000"),
        ("💧 Água",        "4.000.000 m³"),
        ("👷 Mão de Obra", "60.000 hh"),
    ]
    for i, (nome, val) in enumerate(recursos_info):
        top = Inches(1.75) + Inches(0.8) * i
        _rect(s, col2_x, top, col2_w, Inches(0.72),
              CINZA_ESCURO if i % 2 == 0 else RGBColor(0x0F, 0x22, 0x14))
        _txt(s, col2_x + Inches(0.1), top + Inches(0.05),
             Inches(2.8), Inches(0.35), nome, size=12, color=VERDE_CLARO)
        _txt(s, col2_x + Inches(0.1), top + Inches(0.35),
             col2_w - Inches(0.2), Inches(0.35), val, size=16, bold=True, color=BRANCO)

    _rect(s, col2_x, Inches(5.0), col2_w, Inches(0.5), VERDE)
    _txt(s, col2_x + Inches(0.1), Inches(5.05), col2_w, Inches(0.4),
         "Objetivo", size=14, bold=True, color=AMARELO)
    _txt(s, col2_x + Inches(0.1), Inches(5.55), col2_w, Inches(0.8),
         "Maximizar Z (lucro total em R$) determinando\na alocação ótima de hectares.",
         size=13, color=BRANCO)

    # ── S4 Dados das Culturas ──────────────────────────────────────────
    s = _slide(prs)
    _header(s, "Dados das Culturas — Parâmetros do Modelo",
            "Coeficientes utilizados na Programação Linear")

    headers = ["Cultura", "Margem\n(R$/ha)", "Custo Insumos\n(R$/ha)",
               "Água\n(m³/ha)", "MO\n(hh/ha)", "Demanda\nMáx (ha)"]
    col_ws = [Inches(2.5), Inches(1.7), Inches(2.2), Inches(1.5), Inches(1.5), Inches(1.8)]
    col_xs = [M]
    for w in col_ws[:-1]:
        col_xs.append(col_xs[-1] + w)

    row_h = Inches(0.58)
    top0  = Inches(1.25)

    for j, (hdr, cx, cw) in enumerate(zip(headers, col_xs, col_ws)):
        _rect(s, cx, top0, cw - Inches(0.04), row_h, VERDE)
        _txt(s, cx + Inches(0.06), top0 + Inches(0.05), cw - Inches(0.1), row_h,
             hdr, size=12, bold=True, color=AMARELO)

    rows = [
        [culturas[i],
         f"R$ {br(params['margem'][i], 0)}",
         f"R$ {br(params['custo'][i], 0)}",
         f"{br(params['agua'][i], 0)}",
         f"{params['mao_obra'][i]}",
         f"{br(params['demanda_max'][i], 0)}"]
        for i in range(4)
    ]
    bgs = [CINZA_ESCURO, RGBColor(0x0F, 0x22, 0x14)] * 2
    for i, row in enumerate(rows):
        for j, (cell, cx, cw) in enumerate(zip(row, col_xs, col_ws)):
            _rect(s, cx, top0 + row_h * (i + 1), cw - Inches(0.04), row_h, bgs[i])
            bold = j == 0
            clr  = VERDE_CLARO if j == 1 else BRANCO
            _txt(s, cx + Inches(0.06), top0 + row_h * (i + 1) + Inches(0.1),
                 cw - Inches(0.1), row_h, cell, size=14, bold=bold, color=clr)

    _rect(s, 0, Inches(6.35), W, Inches(0.6), CINZA_ESCURO)
    _txt(s, M, Inches(6.42), CW, Inches(0.5),
         "Maior margem: Algodão (R$ 4.200/ha)  |  Menor custo: Cana (R$ 2.800/ha)  |  "
         "Maior consumo de água: Cana (1.500 m³/ha)",
         size=13, color=AMARELO, align=PP_ALIGN.CENTER)

    # ── S5 Modelagem – Variáveis e FO ─────────────────────────────────
    s = _slide(prs)
    _header(s, "Etapa II — Modelagem Matemática",
            "Variáveis de Decisão e Função Objetivo")

    _txt(s, M, Inches(1.3), CW, Inches(0.4),
         "Variáveis de Decisão", size=18, bold=True, color=AMARELO)
    vars_lines = [
        "x₁  =  hectares plantados de Soja",
        "x₂  =  hectares plantados de Milho",
        "x₃  =  hectares plantados de Algodão",
        "x₄  =  hectares plantados de Cana-de-Açúcar",
        "onde  xᵢ ≥ 0  para todo i ∈ {1, 2, 3, 4}",
    ]
    for i, line in enumerate(vars_lines):
        clr = VERDE_CLARO if i == 4 else BRANCO
        _txt(s, Inches(1.0), Inches(1.75) + Inches(0.5) * i,
             Inches(11.5), Inches(0.45), line, size=16, color=clr,
             italic=(i == 4))

    _rect(s, M, Inches(4.4), CW, Inches(0.04), VERDE_MED)

    _txt(s, M, Inches(4.55), CW, Inches(0.4),
         "Função Objetivo", size=18, bold=True, color=AMARELO)

    _rect(s, M, Inches(5.0), CW, Inches(1.2), CINZA_ESCURO)
    _txt(s, M + Inches(0.2), Inches(5.1), CW - Inches(0.3), Inches(0.5),
         "Maximizar  Z  =  3.500 x₁  +  2.300 x₂  +  4.200 x₃  +  3.400 x₄",
         size=22, bold=True, color=VERDE_CLARO, align=PP_ALIGN.CENTER)
    _txt(s, M + Inches(0.2), Inches(5.65), CW - Inches(0.3), Inches(0.4),
         "coeficientes = margem líquida (R$/ha) de cada cultura",
         size=13, color=BRANCO, align=PP_ALIGN.CENTER, italic=True)

    # ── S6 Restrições ──────────────────────────────────────────────────
    s = _slide(prs)
    _header(s, "Etapa II — Restrições do Modelo",
            "Limitações físicas e de mercado da Cooperativa AgroPrime")

    restricoes = [
        ("R1 — Terra (ha)",         "x₁ + x₂ + x₃ + x₄",
         "≤  5.000",  "Área total cultivável disponível"),
        ("R2 — Orçamento (R$)",     "4.500x₁ + 3.200x₂ + 6.000x₃ + 2.800x₄",
         "≤  20.000.000", "Budget de insumos: sementes, fertilizantes, defensivos"),
        ("R3 — Água (m³)",          "600x₁ + 800x₂ + 500x₃ + 1.500x₄",
         "≤  4.000.000",  "Disponibilidade hídrica para irrigação na safra"),
        ("R4 — Mão de Obra (hh)",   "12x₁ + 15x₂ + 22x₃ + 8x₄",
         "≤  60.000",     "Horas-homem disponíveis no período da safra"),
        ("R5 — Demanda Máx (ha)",   "x₁ ≤ 2.500  |  x₂ ≤ 2.000  |  x₃ ≤ 1.500  |  x₄ ≤ 1.000",
         "",              "Limites de absorção do mercado por cultura"),
        ("R6 — Não-negatividade",   "x₁, x₂, x₃, x₄",
         "≥  0",          "Não é possível plantar área negativa"),
    ]

    for i, (label, expr, bound, desc) in enumerate(restricoes):
        top = Inches(1.25) + Inches(0.99) * i
        bg  = CINZA_ESCURO if i % 2 == 0 else RGBColor(0x0F, 0x22, 0x14)
        _rect(s, M, top, CW, Inches(0.88), bg)
        _txt(s, M + Inches(0.1), top + Inches(0.04),
             Inches(2.6), Inches(0.38), label, size=12, bold=True, color=AMARELO)
        expr_full = expr + ("  " + bound if bound else "")
        _txt(s, Inches(3.0), top + Inches(0.04),
             Inches(7.8), Inches(0.4), expr_full, size=13, bold=True, color=VERDE_CLARO)
        _txt(s, M + Inches(0.1), top + Inches(0.46),
             CW - Inches(0.2), Inches(0.38), desc, size=12, color=BRANCO, italic=True)

    # ── S7 Código Python ───────────────────────────────────────────────
    s = _slide(prs)
    _header(s, "Etapa III — Solução Computacional",
            "Python 3 + PuLP (CBC Solver) — Código completo")

    _rect(s, M, Inches(1.25), CW, Inches(5.8), RGBColor(0x0D, 0x1B, 0x10))
    codigo = (
        "from pulp import *\n\n"
        "prob = LpProblem('AgroPrime', LpMaximize)\n"
        "x1 = LpVariable('Soja',    lowBound=0)\n"
        "x2 = LpVariable('Milho',   lowBound=0)\n"
        "x3 = LpVariable('Algodao', lowBound=0)\n"
        "x4 = LpVariable('Cana',    lowBound=0)\n\n"
        "# Função Objetivo\n"
        "prob += 3500*x1 + 2300*x2 + 4200*x3 + 3400*x4\n\n"
        "# Restrições\n"
        "prob += x1 + x2 + x3 + x4          <= 5000        # Terra\n"
        "prob += 4500*x1+3200*x2+6000*x3+2800*x4 <= 20e6   # Orçamento\n"
        "prob += 600*x1+800*x2+500*x3+1500*x4    <= 4e6    # Água\n"
        "prob += 12*x1+15*x2+22*x3+8*x4          <= 60000  # MO\n"
        "prob += x1<=2500; prob += x2<=2000       # Demanda\n"
        "prob += x3<=1500; prob += x4<=1000\n\n"
        "prob.solve(PULP_CBC_CMD(msg=0))\n"
        "print(f'Z* = R$ {br(value(prob.objective), 0)}')"
    )
    _txt(s, M + Inches(0.15), Inches(1.35), CW - Inches(0.25), Inches(5.6),
         codigo, size=12, color=VERDE_CLARO)

    # ── S8 Cálculo Manual ─────────────────────────────────────────────
    s = _slide(prs)
    _header(s, "Etapa IV — Resolução Manual",
            "Identificar restrições ativas → montar sistema 2×2 → resolver por substituição")

    # Passo 1 — variáveis fixadas
    _txt(s, M, Inches(1.3), CW, Inches(0.38),
         "Passo 1 — Restrições ativas identificadas", size=15, bold=True, color=AMARELO)
    ativas = [
        "R2 Orçamento: 100% consumido  →  restrição ativa",
        "R4 Mão de Obra: 100% consumida  →  restrição ativa",
        "R5 Demanda Soja: x₁ = 2.500 ha  →  no limite",
        "R5 Demanda Cana: x₄ = 1.000 ha  →  no limite",
    ]
    for i, txt in enumerate(ativas):
        _txt(s, M + Inches(0.3), Inches(1.72) + Inches(0.32)*i,
             CW, Inches(0.3), f"• {txt}", size=12, color=BRANCO)

    # Passo 2 — sistema reduzido
    _txt(s, M, Inches(3.05), CW, Inches(0.38),
         "Passo 2 — Sistema reduzido (x₁ = 2.500 e x₄ = 1.000 fixados)", size=15, bold=True, color=AMARELO)
    _rect(s, M, Inches(3.45), CW, Inches(0.85), CINZA_ESCURO)
    _txt(s, M + Inches(0.2), Inches(3.52), CW - Inches(0.3), Inches(0.36),
         "R2:  3.200 x₂ + 6.000 x₃  =  20.000.000 − 11.250.000 − 2.800.000  =  5.950.000  …(I)",
         size=13, color=VERDE_CLARO)
    _txt(s, M + Inches(0.2), Inches(3.88), CW - Inches(0.3), Inches(0.36),
         "R4:       15 x₂ +    22 x₃  =  60.000 − 30.000 − 8.000  =  22.000  …(II)",
         size=13, color=VERDE_CLARO)

    # Passo 3 — resolução
    _txt(s, M, Inches(4.45), CW, Inches(0.38),
         "Passo 3 — Resolução por substituição", size=15, bold=True, color=AMARELO)
    _rect(s, M, Inches(4.85), CW, Inches(1.55), CINZA_ESCURO)
    passos = [
        "De (II):  x₂ = (22.000 − 22 x₃) / 15  …(III)",
        "Subst. (III) em (I):  4.693.333,33 − 4.693,33 x₃ + 6.000 x₃  =  5.950.000",
        "                      1.306,67 x₃  =  1.256.666,67   →   x₃  ≈  961,73 ha  (Algodão)  ✓",
        "Subst. x₃ em (III):   x₂  =  (22.000 − 21.158,06) / 15  =  841,94 / 15  ≈  56,12 ha  (Milho)  ✓",
    ]
    for i, txt in enumerate(passos):
        _txt(s, M + Inches(0.2), Inches(4.92) + Inches(0.35)*i,
             CW - Inches(0.3), Inches(0.32),
             txt, size=12, color=BRANCO if i < 2 else AMARELO)

    # Resultado final
    _rect(s, M, Inches(6.5), CW, Inches(0.68), VERDE)
    _txt(s, M + Inches(0.2), Inches(6.56), CW, Inches(0.55),
         f"Z* = 3.500×2.500 + 2.300×56,12 + 4.200×961,73 + 3.400×1.000  ≈  R$ {br(lucro, 0)}",
         size=15, bold=True, color=AMARELO, align=PP_ALIGN.CENTER)

    # ── S9 Cálculos ────────────────────────────────────────────────────
    s = _slide(prs)
    _header(s, "Cálculos — Verificação da Solução Ótima",
            "Substituindo x* na Função Objetivo e nas Restrições")

    _txt(s, M, Inches(1.3), CW, Inches(0.4),
         "Solução encontrada pelo solver:", size=15, bold=True, color=AMARELO)
    sol_txt = (f"x₁* = {br(qtds[0])} ha (Soja)    "
               f"x₂* = {br(qtds[1])} ha (Milho)    "
               f"x₃* = {br(qtds[2])} ha (Algodão)    "
               f"x₄* = {br(qtds[3])} ha (Cana)")
    _txt(s, M, Inches(1.75), CW, Inches(0.45), sol_txt,
         size=14, color=VERDE_CLARO, bold=True)

    _txt(s, M, Inches(2.35), CW, Inches(0.4),
         "Cálculo da Função Objetivo:", size=15, bold=True, color=AMARELO)
    _rect(s, M, Inches(2.8), CW, Inches(2.2), CINZA_ESCURO)
    calculo = (
        f"Z* = 3.500 × {br(qtds[0], 4)}  +  2.300 × {br(qtds[1], 4)}  +  4.200 × {br(qtds[2], 4)}  +  3.400 × {br(qtds[3], 4)}\n\n"
        f"Z* = {br(receitas[0], 0)}  +  {br(receitas[1], 0)}  +  {br(receitas[2], 0)}  +  {br(receitas[3], 0)}\n\n"
        f"Z* =  R$ {br(lucro, 0)}"
    )
    _txt(s, M + Inches(0.2), Inches(2.9), CW - Inches(0.3), Inches(2.0),
         calculo, size=15, color=BRANCO, align=PP_ALIGN.CENTER)

    _txt(s, M, Inches(5.1), CW, Inches(0.4),
         "Verificação das Restrições Ativas:", size=15, bold=True, color=AMARELO)
    u = uso
    verif = [
        (f"Orçamento:  4500×{br(qtds[0], 0)} + 3200×{br(qtds[1], 0)} + 6000×{br(qtds[2], 0)} + 2800×{br(qtds[3], 0)}",
         f"= R$ {br(u['orcamento'][0], 0)}  ≤  R$ 20.000.000  ✓  (ATIVA — folga zero)",
         True),
        (f"Mão de Obra:  12×{br(qtds[0], 0)} + 15×{br(qtds[1], 0)} + 22×{br(qtds[2], 0)} + 8×{br(qtds[3], 0)}",
         f"= {br(u['mao_obra'][0], 0)} hh  ≤  60.000 hh  ✓  (ATIVA — folga zero)",
         True),
        (f"Terra:  {br(sum(qtds))} ha  ≤  5.000 ha  ✓  (folga: {br(5000-sum(qtds), 0)} ha)",
         "", False),
    ]
    for i, (expr, result, ativa) in enumerate(verif):
        top = Inches(5.55) + Inches(0.6) * i
        clr = VERMELHO if ativa else VERDE_CLARO
        _txt(s, M, top, Inches(7.5), Inches(0.32), expr, size=11, color=BRANCO)
        if result:
            _txt(s, M, top + Inches(0.28), Inches(12.0), Inches(0.28),
                 result, size=11, bold=True, color=clr)

    # ── S9 Gráfico Hectares ────────────────────────────────────────────
    s = _slide(prs)
    _header(s, "Resultados — Alocação Ótima de Hectares",
            f"Lucro Total Máximo: R$ {br(lucro, 0)}")

    fig = _grafico_barras(culturas, qtds, "Hectares por Cultura (Solução Ótima)")
    buf = _fig_to_bytes(fig)
    _add_img(s, buf, M, Inches(1.3), Inches(7.2), Inches(4.5))
    plt.close(fig)

    _rect(s, Inches(8.0), Inches(1.3), Inches(4.8), Inches(4.5), CINZA_ESCURO)
    _txt(s, Inches(8.1), Inches(1.4), Inches(4.6), Inches(0.5),
         "Participação na Receita", size=14, bold=True, color=AMARELO)
    for i, (c, q, r) in enumerate(zip(culturas, qtds, receitas)):
        top = Inches(1.95) + Inches(0.72) * i
        _rect(s, Inches(8.1), top, Inches(0.18), Inches(0.4),
              RGBColor(*[int(CORES_CULT[i].lstrip("#")[j:j+2], 16) for j in (0, 2, 4)]))
        _txt(s, Inches(8.35), top, Inches(2.2), Inches(0.4),
             c, size=12, color=BRANCO)
        _txt(s, Inches(8.35), top + Inches(0.22), Inches(4.4), Inches(0.38),
             f"{br(q, 0)} ha  →  R$ {br(r, 0)}", size=11, color=VERDE_CLARO)

    _rect(s, M, Inches(5.95), CW, Inches(0.7), VERDE)
    _txt(s, M + Inches(0.2), Inches(6.0), CW, Inches(0.6),
         f"LUCRO ÓTIMO TOTAL:  R$ {br(lucro, 0)}", size=22, bold=True,
         color=AMARELO, align=PP_ALIGN.CENTER)

    # ── S10 Uso dos Recursos ───────────────────────────────────────────
    s = _slide(prs)
    _header(s, "Resultados — Uso dos Recursos",
            "Restrições ativas identificam os gargalos da operação")

    fig2 = _grafico_recursos(uso)
    buf2 = _fig_to_bytes(fig2)
    _add_img(s, buf2, M, Inches(1.3), Inches(8.0), Inches(4.0))
    plt.close(fig2)

    labels_rec = {"terra": "Terra (ha)", "orcamento": "Orçamento (R$)",
                  "agua": "Água (m³)", "mao_obra": "Mão de Obra (hh)"}
    _rect(s, Inches(8.8), Inches(1.3), Inches(4.0), Inches(4.0), CINZA_ESCURO)
    _txt(s, Inches(8.9), Inches(1.4), Inches(3.8), Inches(0.45),
         "Resumo dos Recursos", size=14, bold=True, color=AMARELO)
    for i, (k, (us, tot, un)) in enumerate(uso.items()):
        pct = us / tot * 100
        top = Inches(1.92) + Inches(0.77) * i
        clr = VERMELHO if pct >= 99 else VERDE_CLARO
        _txt(s, Inches(8.9), top, Inches(3.8), Inches(0.3),
             labels_rec[k], size=11, bold=True, color=AMARELO)
        _txt(s, Inches(8.9), top + Inches(0.28), Inches(3.8), Inches(0.28),
             f"{br(us, 0)} / {br(tot, 0)} {un}  ({pct:.1f}%)",
             size=11, color=clr)

    _rect(s, M, Inches(5.5), CW, Inches(0.04), VERDE_MED)
    _txt(s, M, Inches(5.6), CW, Inches(0.7),
         "Gargalos identificados: Orçamento e Mão de Obra estão 100% utilizados. "
         "Investir em mais crédito rural ou mecanização liberaria capacidade adicional de produção.",
         size=13, color=BRANCO)

    # ── S11 Método Gráfico ─────────────────────────────────────────────
    s = _slide(prs)
    _header(s, "Método Gráfico — Região Viável e Ponto Ótimo",
            "Projeção 2D: Soja (x₁) × Algodão (x₃) — Milho e Cana fixos nos valores ótimos")

    x2_fix, x4_fix = qtds[1], qtds[3]
    terra_adj = 5000       - x2_fix        - x4_fix
    orc_adj   = 20_000_000 - 3200*x2_fix   - 2800*x4_fix
    mo_adj    = 60_000     - 15*x2_fix     - 8*x4_fix
    dem1, dem3 = 2500.0, 1500.0

    x_max, y_max = 2700.0, 1600.0
    xv = np.linspace(0, x_max, 1000)

    fig_g, ax_g = plt.subplots(figsize=(8, 5), facecolor="#F0F4F0")
    ax_g.set_facecolor("#FFFFFF")

    ax_g.plot(xv, terra_adj - xv,              color="#27AE60", lw=2.2, label="Terra")
    ax_g.plot(xv, (orc_adj - 4500*xv) / 6000,  color="#E67E22", lw=2.2, label="Orçamento")
    ax_g.plot(xv, (mo_adj - 12*xv) / 22,        color="#C0392B", lw=2.5, label="Mão de Obra")
    ax_g.axvline(dem1, color="#2980B9", lw=2.0, ls="--", label=f"Dem. Soja ≤ {dem1:.0f}")
    ax_g.axhline(dem3, color="#8E44AD", lw=2.0, ls="--", label=f"Dem. Algodão ≤ {dem3:.0f}")

    xs = np.linspace(0, x_max, 3000)
    y_up = np.minimum(np.minimum(terra_adj - xs, (orc_adj - 4500*xs)/6000),
                      np.minimum((mo_adj - 12*xs)/22, dem3))
    y_up = np.clip(y_up, 0, y_max)
    mask = (xs >= 0) & (xs <= dem1)
    ax_g.fill_between(xs[mask], 0, y_up[mask], alpha=0.20, color="#27AE60", label="Região viável")

    Z_opt = 3500*qtds[0] + 4200*qtds[2]
    for frac in [0.5, 0.75, 1.0]:
        Z = Z_opt * frac
        y_iso = (Z - 3500*xv) / 4200
        ax_g.plot(xv, y_iso, color="#95A5A6", lw=2.0 if frac==1.0 else 1.0,
                  ls="-" if frac==1.0 else ":", alpha=0.9 if frac==1.0 else 0.6,
                  label=f"Isolucro R${Z/1e6:.1f}M" if frac==1.0 else None)

    eps = 0.5
    verts = []
    try:
        p_ = np.linalg.solve([[4500,6000],[12,22]], [orc_adj, mo_adj])
        if -eps<=p_[0]<=dem1+eps and -eps<=p_[1]<=dem3+eps:
            verts.append((p_[0], p_[1], "Orç∩MO"))
    except Exception:
        pass
    x3_d1 = (mo_adj - 12*dem1)/22
    if -eps<=x3_d1<=dem3+eps:
        verts.append((dem1, x3_d1, "Ótimo"))
    x1_d3 = (mo_adj - 22*dem3)/12
    if -eps<=x1_d3<=dem1+eps:
        verts.append((x1_d3, dem3, "MO∩Dem.Algodão"))

    for vx, vy, vlabel in verts:
        is_opt = abs(vx-qtds[0])<2 and abs(vy-qtds[2])<2
        c = "#F39C12" if is_opt else "#2C3E50"
        ax_g.scatter(vx, vy, color=c, s=120 if is_opt else 60, zorder=6, edgecolors="white", linewidth=1.5)
        ax_g.annotate(f"({vx:.0f}, {vy:.0f})", xy=(vx,vy),
                      xytext=(vx-400 if is_opt else vx+40, vy+100),
                      color=c, fontsize=9, fontweight="bold",
                      arrowprops=dict(arrowstyle="->", color=c, lw=1.2))

    ax_g.scatter(qtds[0], qtds[2], color="#F39C12", s=220, zorder=10,
                 marker="*", edgecolors="#1A2E1D", linewidth=1)
    ax_g.annotate(f"★ ÓTIMO\nSoja={qtds[0]:.0f} ha\nAlgodão={qtds[2]:.0f} ha",
                  xy=(qtds[0],qtds[2]), xytext=(qtds[0]-850, qtds[2]+250),
                  color="#C0392B", fontsize=10, fontweight="bold",
                  arrowprops=dict(arrowstyle="->", color="#C0392B", lw=1.5))

    ax_g.set_xlim(0, x_max)
    ax_g.set_ylim(0, y_max)
    ax_g.set_xlabel("x₁ — Soja (ha)", color="#1A2E1D", fontsize=11, fontweight="bold")
    ax_g.set_ylabel("x₃ — Algodão (ha)", color="#1A2E1D", fontsize=11, fontweight="bold")
    ax_g.set_title("Região Viável e Ponto Ótimo (Simplex)", color="#1A2E1D", fontsize=12, fontweight="bold", pad=10)
    ax_g.tick_params(colors="#1A2E1D", labelsize=9)
    for sp in ax_g.spines.values():
        sp.set_edgecolor("#AAAAAA")
    ax_g.legend(loc="lower center", bbox_to_anchor=(0.5, -0.18),
                ncols=3, fontsize=8, facecolor="#F0F4F0",
                edgecolor="#AAAAAA", labelcolor="#1A2E1D")
    ax_g.xaxis.grid(True, color="#DDDDDD", linewidth=0.6)
    ax_g.yaxis.grid(True, color="#DDDDDD", linewidth=0.6)
    ax_g.set_axisbelow(True)
    plt.subplots_adjust(bottom=0.22)

    buf_g = _fig_to_bytes(fig_g)
    _add_img(s, buf_g, M, Inches(1.25), Inches(8.5), Inches(5.5))
    plt.close(fig_g)

    _rect(s, Inches(9.2), Inches(1.4), Inches(3.6), Inches(5.4), CINZA_ESCURO)
    _txt(s, Inches(9.3), Inches(1.5), Inches(3.4), Inches(0.45),
         "Como ler o gráfico", size=13, bold=True, color=AMARELO)
    notas = [
        ("Retas coloridas", "Cada linha = uma restrição. A solução deve ficar abaixo/esquerda de todas."),
        ("Região viável", "Área verde: todos os pontos que satisfazem as restrições simultaneamente."),
        ("Pontos de vértice", "Em LP o ótimo SEMPRE fica em um vértice (canto) da região viável."),
        ("Isolucros (cinza)", "Combinações com mesmo lucro. Quanto mais alto/direita, maior o lucro."),
        ("★ Ponto ótimo", "Onde a isolucro mais alta toca a fronteira da região viável."),
    ]
    for i, (titulo_nota, desc_nota) in enumerate(notas):
        top = Inches(2.05) + Inches(0.95) * i
        _txt(s, Inches(9.3), top, Inches(3.4), Inches(0.3),
             titulo_nota, size=10, bold=True, color=VERDE_CLARO)
        _txt(s, Inches(9.3), top + Inches(0.28), Inches(3.4), Inches(0.55),
             desc_nota, size=9, color=BRANCO)

    # ── S12 Conclusão ──────────────────────────────────────────────────
    s = _slide(prs, VERDE)
    _rect(s, 0, Inches(2.3), W, Inches(3.2), PRETO_VERDE)
    _rect(s, 0, Inches(5.5), W, Inches(0.06), AMARELO)

    _txt(s, M, Inches(0.3), CW, Inches(0.6),
         "CONCLUSÃO", size=16, bold=True, color=AMARELO, align=PP_ALIGN.CENTER)
    _txt(s, M, Inches(0.9), CW, Inches(1.1),
         "A Programação Linear encontrou o mix ótimo de culturas para a Cooperativa AgroPrime",
         size=20, color=BRANCO, align=PP_ALIGN.CENTER)

    _rect(s, Inches(2.0), Inches(2.4), Inches(9.33), Inches(1.2), CINZA_ESCURO)
    _txt(s, Inches(2.1), Inches(2.5), Inches(9.0), Inches(1.0),
         f"Lucro Máximo:  R$ {br(lucro, 0)}",
         size=36, bold=True, color=AMARELO, align=PP_ALIGN.CENTER)

    _txt(s, M, Inches(3.75), CW, Inches(0.5),
         f"Soja: {br(qtds[0], 0)} ha  |  Algodão: {br(qtds[2], 0)} ha  |  Cana: {br(qtds[3], 0)} ha  |  Milho: {br(qtds[1], 0)} ha",
         size=16, color=VERDE_CLARO, align=PP_ALIGN.CENTER)

    insights = [
        "Soja e Cana dominam a alocação — Soja pelo volume (limite de demanda), Cana pela melhor eficiência de custo e mão de obra",
        "Orçamento e Mão de Obra são os gargalos — ampliar esses recursos aumenta o lucro",
        "O modelo é replicável para qualquer cooperativa agrícola real",
    ]
    for i, ins in enumerate(insights):
        _txt(s, M, Inches(4.45) + Inches(0.42) * i, CW, Inches(0.38),
             f"• {ins}", size=13, color=BRANCO, align=PP_ALIGN.CENTER)

    _txt(s, M, Inches(6.5), CW, Inches(0.5),
         "MBA em Pesquisa Operacional e Tomada de Decisão  |  BSBr  |  Prof. Gabriel Capela  |  Jun/2026",
         size=12, color=BRANCO, align=PP_ALIGN.CENTER)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
